"""Structural QA for PPTX — checks positioning, overflow, margins, overlaps."""
import sys
from pptx import Presentation
from pptx.util import Inches, Emu

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MIN_MARGIN = Inches(0.4)

def emu_to_inches(emu):
    return emu / 914400

def check_slide(slide, idx):
    issues = []
    shapes = list(slide.shapes)

    for shape in shapes:
        x = emu_to_inches(shape.left or 0)
        y = emu_to_inches(shape.top or 0)
        w = emu_to_inches(shape.width or 0)
        h = emu_to_inches(shape.height or 0)
        r = x + w
        b = y + h

        name = f"{shape.shape_type} '{shape.name}'"

        # Check bounds
        if x < -0.1:
            issues.append(f"  LEFT OVERFLOW: {name} x={x:.2f}")
        if r > 10.1:
            issues.append(f"  RIGHT OVERFLOW: {name} right={r:.2f}")
        if b > 5.7:
            issues.append(f"  BOTTOM OVERFLOW: {name} bottom={b:.2f}")

        # Check margin (skip background shapes)
        is_bg = (w > Inches(9) and h > Inches(5)) or shape.name.startswith("Background")
        if not is_bg and x > 0 and x < MIN_MARGIN:
            issues.append(f"  MARGIN: {name} left={x:.2f}in < {emu_to_inches(MIN_MARGIN):.1f}in")

        # Check text overflow
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text.strip()
            if text and w > 0:
                # Rough estimate: Chinese chars ~12pt wide at 12pt font
                max_chars_per_line = max(1, int(w / 0.15))
                longest = max((len(p.text) for p in tf.paragraphs), default=0)
                if longest > max_chars_per_line * 3 and h < Inches(0.5):
                    issues.append(f"  TEXT MAY OVERFLOW: {name} max_line={longest}chars, box_w={w:.2f}in, box_h={h:.2f}in")

    # Check overlaps between text-bearing shapes
    for i, a in enumerate(shapes):
        if not a.has_text_frame or not a.text_frame.text.strip():
            continue
        ax = emu_to_inches(a.left or 0)
        ay = emu_to_inches(a.top or 0)
        aw = emu_to_inches(a.width or 0)
        ah = emu_to_inches(a.height or 0)

        for b in shapes[i+1:]:
            if not b.has_text_frame or not b.text_frame.text.strip():
                continue
            bx = emu_to_inches(b.left or 0)
            by = emu_to_inches(b.top or 0)
            bw = emu_to_inches(b.width or 0)
            bh = emu_to_inches(b.height or 0)

            # Check overlap
            if (ax < bx + bw and ax + aw > bx and ay < by + bh and ay + ah > by):
                overlap_x = min(ax+aw, bx+bw) - max(ax, bx)
                overlap_y = min(ay+ah, by+bh) - max(ay, by)
                if overlap_x > 0.1 and overlap_y > 0.1:
                    a_text = a.text_frame.text[:30]
                    b_text = b.text_frame.text[:30]
                    issues.append(f"  OVERLAP: '{a_text}' ↔ '{b_text}' area=({overlap_x:.1f}x{overlap_y:.1f}in)")

    if issues:
        print(f"Slide {idx+1} ISSUES:")
        for issue in issues:
            print(issue)
    else:
        print(f"Slide {idx+1}: OK")
    print()

prs = Presentation("output/UXB项目团队分享.pptx")
print(f"Slides: {len(prs.slides)}, Size: {prs.slide_width}x{prs.slide_height}\n")

for i, slide in enumerate(prs.slides):
    check_slide(slide, i)
